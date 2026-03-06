#!/usr/bin/env python3
"""Export HTML slide presentation to PowerPoint."""

import os
import urllib.parse
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches

# Paths
HTML_PATH = Path("/Users/ericcoffie/Encore Funding/tempnet-staffing-presentation-v4.html")
OUTPUT_PATH = Path("/Users/ericcoffie/Downloads/TempNet_Presentation_v4.pptx")
TEMP_DIR = Path("/Users/ericcoffie/Encore Funding/.slide_screenshots")

# Slide dimensions (16:9 aspect ratio matching 960x540)
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)

def main():
    # Create temp directory
    TEMP_DIR.mkdir(exist_ok=True)

    file_url = "file://" + urllib.parse.quote(str(HTML_PATH))

    print("Launching browser...")
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": 960, "height": 540})

        print(f"Loading {HTML_PATH.name}...")
        page.goto(file_url, wait_until="networkidle")

        # Get all slides
        slides = page.query_selector_all(".slide")
        total = len(slides)
        print(f"Found {total} slides")

        # Screenshot each slide
        screenshot_paths = []
        for i, slide in enumerate(slides):
            print(f"  Capturing slide {i + 1}/{total}...")
            screenshot_path = TEMP_DIR / f"slide_{i + 1:02d}.png"
            slide.screenshot(path=str(screenshot_path))
            screenshot_paths.append(screenshot_path)

        browser.close()

    # Create PowerPoint
    print("Creating PowerPoint...")
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Blank layout
    blank_layout = prs.slide_layouts[6]

    for i, img_path in enumerate(screenshot_paths):
        print(f"  Adding slide {i + 1}/{total}...")
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path),
            Inches(0), Inches(0),
            width=SLIDE_WIDTH,
            height=SLIDE_HEIGHT
        )

    # Save
    prs.save(str(OUTPUT_PATH))
    print(f"\nSaved: {OUTPUT_PATH}")

    # Cleanup screenshots
    for p in screenshot_paths:
        p.unlink()
    TEMP_DIR.rmdir()
    print("Done!")

if __name__ == "__main__":
    main()
